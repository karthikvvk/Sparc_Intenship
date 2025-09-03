import requests
from tempfile import NamedTemporaryFile
from pathlib import Path
import mimetypes

def extract_text_from_url(file_url: str):
    print(file_url)
    try:
        # ---------------------- Handle local file ---------------------- #
        if Path(file_url).exists():
            ext = Path(file_url).suffix.lower().lstrip(".")
            tmp_path = Path(file_url)

        # ---------------------- Handle remote URL ---------------------- #
        elif file_url.startswith("http://") or file_url.startswith("https://"):
            head = requests.head(file_url, allow_redirects=True, timeout=10)
            size = int(head.headers.get("Content-Length", 0) or 0)
            if size and size > 500 * 1024 * 1024:
                return "⚠ File too large to process.", None

            ext = file_url.split('?')[0].split('.')[-1].lower()
            if not ext:
                ext = (mimetypes.guess_extension(head.headers.get("Content-Type", "")) or "bin").lstrip('.')

            r = requests.get(file_url, timeout=30)
            r.raise_for_status()

            with NamedTemporaryFile(delete=False, suffix=f".{ext}") as f:
                f.write(r.content)
                f.flush()
                tmp_path = Path(f.name)
        else:
            return "⚠ Invalid file path or URL", None

        # ---------------------- PDF Handling ---------------------- #
        if ext == "pdf":
            import fitz  # PyMuPDF
            from PIL import Image
            import pytesseract
            import io

            doc = fitz.open(str(tmp_path))
            pdf_text = []
            for page_num, page in enumerate(doc):
                text = page.get_text()
                if text.strip():
                    pdf_text.append(text)
                else:
                    # Fallback to OCR if page is image-only
                    pix = page.get_pixmap(dpi=300)
                    img = Image.open(io.BytesIO(pix.tobytes("png")))
                    ocr_text = pytesseract.image_to_string(img)
                    if ocr_text.strip():
                        pdf_text.append(ocr_text)
            return ("\n".join(pdf_text) if pdf_text else "⚠ No readable text found in PDF."), ext

        # ---------------------- DOCX Handling ---------------------- #
        if ext == "docx":
            import docx
            doc = docx.Document(str(tmp_path))
            return "\n".join(p.text for p in doc.paragraphs), ext

        # ---------------------- EML Handling ---------------------- #
        if ext == "eml":
            from bs4 import BeautifulSoup
            html = tmp_path.read_text(errors="ignore")
            return BeautifulSoup(html, "html.parser").get_text("\n"), ext

        # ---------------------- PPTX Handling ---------------------- #
        if ext == "pptx":
            from pptx import Presentation
            from PIL import Image
            import pytesseract, io
            prs = Presentation(str(tmp_path))
            slides = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    try:
                        if hasattr(shape, "image"):
                            img = Image.open(io.BytesIO(shape.image.blob))
                            ocr = pytesseract.image_to_string(img)
                            if ocr.strip():
                                slide_text.append(ocr.strip())
                    except Exception:
                        pass
                if slide_text:
                    slides.append("\n".join(slide_text))
            return ("\n".join(slides) if slides else "⚠ No readable text found in PPTX."), ext

        # ---------------------- Image Handling ---------------------- #
        if ext in ("png", "jpg", "jpeg"):
            from PIL import Image
            import pytesseract
            img = Image.open(str(tmp_path))
            txt = pytesseract.image_to_string(img)
            return (txt if txt.strip() else "⚠ No readable text found in image."), ext

        # ---------------------- Excel Handling ---------------------- #
        if ext in ("xlsx", "xls"):
            import openpyxl
            wb = openpyxl.load_workbook(str(tmp_path), data_only=True)
            rows = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    rows.append(" ".join(str(c) if c is not None else "" for c in row))
            return ("\n".join(rows) if rows else "⚠ No text found in XLSX."), ext

        return f"⚠ Unsupported file format: {ext}", ext

    except Exception as e:
        return f"⚠ Error processing file: {e}", None
